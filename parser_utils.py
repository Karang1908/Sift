import os
import time
import logging
import subprocess
import pypdf
import docx
import openpyxl
import xlrd
import pytesseract
from PIL import Image
from pptx import Presentation

logger = logging.getLogger("docparser.parser")

# Legacy formats textutil (built into macOS) can convert to plain text.
LEGACY_TEXTUTIL_EXTENSIONS = {'.doc', '.rtf', '.odt', '.wordml', '.webarchive'}
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp'}


def _convert_with_textutil(filepath: str) -> str:
    result = subprocess.run(
        ["textutil", "-convert", "txt", "-stdout", filepath],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr.strip() or "textutil conversion failed")
    return result.stdout


def _ocr_image(filepath: str) -> str:
    with Image.open(filepath) as image:
        return pytesseract.image_to_string(image)


def _extract_openpyxl_workbook(wb) -> list:
    lines = []
    try:
        for sheet_name in wb.sheetnames:
            lines.append(f"\n--- [Sheet: {sheet_name}] ---")
            sheet = wb[sheet_name]
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True)):
                if any(val is not None and str(val).strip() for val in row):
                    row_str = " | ".join([str(val) if val is not None else "" for val in row])
                    lines.append(f"[Row {row_idx+1}] {row_str}")
    finally:
        wb.close()
    return lines


# Every failure/unsupported branch below appends a bracketed sentinel string
# as the ENTIRE content instead of raising, so callers (app.py's upload/list
# endpoints) can use this to tell "parsed successfully" apart from "the cache
# file exists but only contains a placeholder" - otherwise a failed parse
# silently looks identical to a successful one.
_ERROR_SENTINEL_PREFIXES = ("[Error ", "[Unsupported", "[Legacy", "[Image contains no recognizable text")


def is_error_content(text: str) -> bool:
    stripped = text.strip()
    return any(stripped.startswith(prefix) for prefix in _ERROR_SENTINEL_PREFIXES)


def _looks_binary(filepath: str, sample_size: int = 8192) -> bool:
    with open(filepath, 'rb') as f:
        sample = f.read(sample_size)
    if not sample:
        return False
    if b'\x00' in sample:
        return True
    text_bytes = bytes(range(32, 127)) + b'\n\r\t\f\b'
    nontext = sum(b not in text_bytes for b in sample)
    return (nontext / len(sample)) > 0.30


def extract_text_from_file(filepath: str) -> str:
    start = time.perf_counter()
    filename = os.path.basename(filepath)
    _, ext = os.path.splitext(filename.lower())
    content = []
    
    # Text-based formats
    text_extensions = {'.txt', '.md', '.py', '.json', '.csv', '.tsv', '.html', '.xml', '.yaml', '.yml', '.ini', '.conf'}
    
    if ext in text_extensions:
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content.append(f.read())
        except Exception as e:
            logger.warning("error reading text file %s: %s", filename, e)
            content.append(f"[Error reading text file: {str(e)}]")
            
    elif ext == '.pdf':
        try:
            reader = pypdf.PdfReader(filepath)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                content.append(f"\n--- [Page {i+1}] ---\n{page_text}")
        except Exception as e:
            logger.warning("error parsing PDF %s: %s", filename, e)
            content.append(f"[Error parsing PDF: {str(e)}]")
            
    elif ext == '.docx':
        try:
            doc = docx.Document(filepath)
            # Paragraphs
            for para_idx, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    content.append(f"[Paragraph {para_idx+1}] {para.text}")
            # Tables
            for tbl_idx, table in enumerate(doc.tables):
                content.append(f"\n--- [Table {tbl_idx+1}] ---")
                for row_idx, row in enumerate(table.rows):
                    row_text = [cell.text.strip() for cell in row.cells]
                    # Filter out empty cells/rows
                    if any(row_text):
                        content.append(f"[Row {row_idx+1}] " + " | ".join(row_text))
        except Exception as e:
            logger.warning("error parsing Word doc %s: %s", filename, e)
            content.append(f"[Error parsing Word doc: {str(e)}]")
            
    elif ext == '.pptx':
        try:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                content.append(f"\n--- [Slide {i+1}] ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content.append(shape.text)
        except Exception as e:
            logger.warning("error parsing PowerPoint slide %s: %s", filename, e)
            content.append(f"[Error parsing PowerPoint slide: {str(e)}]")

    elif ext == '.ppt':
        # python-pptx cannot read the legacy binary .ppt format, and there's
        # no lightweight pure-Python/macOS-builtin way to convert it - report
        # honestly rather than silently garbage-decoding the binary content.
        logger.warning("legacy .ppt format not supported: %s", filename)
        content.append("[Legacy PowerPoint .ppt format is not supported - please save as .pptx and re-upload]")

    elif ext == '.xlsx':
        try:
            # read_only avoids materializing the whole workbook in memory;
            # must close() explicitly or the underlying zip archive stays open.
            wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
            content.extend(_extract_openpyxl_workbook(wb))
        except Exception as e:
            logger.warning("error parsing Excel sheet %s: %s", filename, e)
            content.append(f"[Error parsing Excel sheet: {str(e)}]")

    elif ext == '.xls':
        try:
            # Some "legacy" .xls files are actually modern .xlsx renamed -
            # try openpyxl first, fall back to xlrd for true legacy BIFF files.
            try:
                wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
                content.extend(_extract_openpyxl_workbook(wb))
            except Exception:
                book = xlrd.open_workbook(filepath)
                for sheet in book.sheets():
                    content.append(f"\n--- [Sheet: {sheet.name}] ---")
                    for row_idx in range(sheet.nrows):
                        row = sheet.row_values(row_idx)
                        if any(str(val).strip() for val in row):
                            row_str = " | ".join(str(val) for val in row)
                            content.append(f"[Row {row_idx+1}] {row_str}")
        except Exception as e:
            logger.warning("error parsing Excel sheet %s: %s", filename, e)
            content.append(f"[Error parsing Excel sheet: {str(e)}]")

    elif ext in LEGACY_TEXTUTIL_EXTENSIONS:
        try:
            content.append(_convert_with_textutil(filepath))
        except Exception as e:
            logger.warning("error converting legacy document %s: %s", filename, e)
            content.append(f"[Error converting document: {str(e)}]")

    elif ext in IMAGE_EXTENSIONS:
        try:
            ocr_text = _ocr_image(filepath)
            content.append(ocr_text if ocr_text.strip() else "[Image contains no recognizable text (OCR found nothing)]")
        except Exception as e:
            logger.warning("error running OCR on image %s: %s", filename, e)
            content.append(f"[Error running OCR on image: {str(e)}]")

    else:
        # Fallback: check whether this actually looks like text before
        # decoding - otherwise arbitrary binary data (zip, exe, audio, ...)
        # would be silently decoded into garbage "text" instead of reporting
        # the format as unsupported.
        if _looks_binary(filepath):
            logger.warning("binary/unsupported file format: %s", filename)
            content.append("[Unsupported binary file format - no text extracted]")
        else:
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content.append(f.read())
            except Exception as e:
                logger.warning("unsupported file format or read error %s: %s", filename, e)
                content.append(f"[Unsupported file format or read error: {str(e)}]")

    result = "\n".join(content)
    logger.info(
        "parsed %s (%s) in %.2fs -> %d chars",
        filename, ext or "no-ext", time.perf_counter() - start, len(result),
    )
    return result
